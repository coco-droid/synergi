import os
import subprocess
from selenium import webdriver
from selenium.webdriver.common.by import By
from selenium.webdriver.common.keys import Keys
from selenium.webdriver.support.ui import WebDriverWait
import requests
from selenium.webdriver.support import expected_conditions as EC
from docx import Document
from docx2pdf import convert
from pptx import Presentation
from openpyxl import Workbook

class Executor:
    def __init__(self, websocket):
        self.websocket = websocket
        self.os = self.detect_os()

    def detect_os(self):
        import platform
        return platform.system()

    def execute(self, action, send_response=False):
        try:
            if action["tool"] == "app":
                return self.launch_app(action["name"])
            elif action["tool"] == "bash":
                return self.bash_execute(action["command"])
            elif action["tool"] == "search":
                return self.web_search(action["query"], action["search_engine"])
            elif action["tool"] == "create_folder":
                return self.create_folder(action["folder_path"])
            elif action["tool"] == "create_file":
                return self.create_file(action["file_path"])
            elif action["tool"] == "write_to_file":
                return self.write_to_file(action["file_path"], action["content"])
            elif action["tool"] == "delete_file":
                return self.delete_file(action["file_path"])
            elif action["tool"] == "copy_file":
                return self.copy_file(action["source_path"], action["destination_path"])
            elif action["tool"] == "move_file":
                return self.move_file(action["source_path"], action["destination_path"])
            elif action["tool"] == "create_word_doc":
                return self.create_word_doc(action["file_path"], action["content"])
            elif action["tool"] == "create_pdf":
                return self.create_pdf(action["word_file_path"], action["pdf_file_path"])
            elif action["tool"] == "convert_latex_to_pdf":
                return self.convert_latex_to_pdf(action["latex_file"])
            elif action["tool"] == "convert_latex_to_docx":
                return self.convert_latex_to_docx(action["latex_file"])
            elif action["tool"] == "convert_markdown_to_pdf":
                return self.convert_markdown_to_pdf(action["markdown_file"])
            elif action["tool"] == "convert_markdown_to_docx":
                return self.convert_markdown_to_docx(action["markdown_file"])
            else:
                raise ValueError("Invalid tool specified.")
        except Exception as e:
            if send_response:
                error_msg = self.generate_error_response(action, e)
                # self.websocket.send(error_msg)

    def launch_app(self, name):
        if self.os == "Windows":
            os.startfile(name)
        else:
            os.system(f"open {name}")

    def bash_execute(self, command):
        return subprocess.check_output(command, shell=True)

    def web_search(self, query, search_engine):
        if search_engine == "google":
            driver = webdriver.Chrome()
            driver.get("https://www.google.com")
            search_box = driver.find_element_by_name("q")
            search_box.send_keys(query)
            search_box.submit()
            return driver.title
        elif search_engine == "duckduckgo":
            driver = webdriver.Chrome()
            driver.get("https://www.duckduckgo.com")
            search_box = driver.find_element_by_name("q")
            search_box.send_keys(query)
            search_box.submit()
            return driver.title
        elif search_engine == "bing":
            driver = webdriver.Chrome()
            driver.get("https://www.bing.com")
            search_box = driver.find_element_by_name("q")
            search_box.send_keys(query)
            search_box.submit()
            return driver.title
        elif search_engine == "phind":
            return results

    def create_folder(self, folder_path):
        os.makedirs(folder_path)

    def create_file(self, file_path):
        with open(file_path, "w") as file:
            pass

    def write_to_file(self, file_path, content):
        with open(file_path, "w") as file:
            file.write(content)

    def delete_file(self, file_path):
        os.remove(file_path)

    def copy_file(self, source_path, destination_path):
        shutil.copy2(source_path, destination_path)

    def move_file(self, source_path, destination_path):
        shutil.move(source_path, destination_path)

    def create_word_doc(self, file_path, content):
        doc = Document()
        doc.add_paragraph(content)
        doc.save(file_path)

    def create_pdf(self, word_file_path, pdf_file_path):
        convert(word_file_path, pdf_file_path)

    def convert_latex_to_pdf(self, latex_file):
        subprocess.run(['pdflatex', '-interaction', 'nonstopmode', latex_file])

    def convert_latex_to_docx(self, latex_file):
        subprocess.run(['pandoc', '-s', latex_file, '-o', 'output.docx'])

    def convert_markdown_to_pdf(self, markdown_file):
        subprocess.run(['markdown2pdf', markdown_file])

    def convert_markdown_to_docx(self, markdown_file):
        subprocess.run(['pandoc', '-s', markdown_file, '-o', 'output.docx'])

    def create_powerpoint_file():
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        title = slide.shapes.title
        title.text = "Titre de la diapositive"
        prs.save("presentation.pptx")
    def create_excel_file():
       wb = Workbook()
       ws = wb.active
       ws['A1'] = "Contenu de la cellule A1"
       wb.save("excel.xlsx")

# Example usage
if __name__ == '__main__':
    create_powerpoint_file()
    create_excel_file()

    # Example usage
    latex_file = 'input.tex'
    markdown_file = 'input.md'

    Executor().convert_latex_to_pdf(latex_file)
    Executor().convert_latex_to_docx(latex_file)
    Executor().convert_markdown_to_pdf(markdown_file)
    Executor().convert_markdown_to_docx(markdown_file)